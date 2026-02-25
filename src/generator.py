"""スライド生成モジュール (Step 4)

data/analyzed_report.json から PowerPoint (.pptx) ファイルを生成し
output/ ディレクトリに保存する。
"""

import json
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from utils import ensure_dirs, setup_logger

logger = setup_logger(__name__)

INPUT_PATH = Path("data/analyzed_report.json")
OUTPUT_DIR = Path("output")

# ---------------------------------------------------------------------------
# カラーパレット（ダークネイビー系）
# ---------------------------------------------------------------------------
COLOR_BG = RGBColor(0x0D, 0x1B, 0x2A)      # 濃紺（背景）
COLOR_TITLE = RGBColor(0xFF, 0xFF, 0xFF)    # 白（タイトル）
COLOR_BODY = RGBColor(0xD0, 0xE4, 0xFF)     # 薄青白（本文）
COLOR_ACCENT = RGBColor(0x4F, 0xB3, 0xBF)  # ティール（アクセント・区切り線）
COLOR_SOURCE = RGBColor(0x80, 0xA0, 0xC0)  # 灰青（出典）
COLOR_NUM = RGBColor(0x4F, 0xB3, 0xBF)     # ティール（スライド番号）


# ---------------------------------------------------------------------------
# ユーティリティ
# ---------------------------------------------------------------------------


def _set_bg(slide, color: RGBColor) -> None:
    """スライド背景色を単色で設定する"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    font_size: int,
    bold: bool = False,
    color: RGBColor = COLOR_BODY,
    align=PP_ALIGN.LEFT,
) -> None:
    """テキストボックスをスライドに追加する"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_divider(slide, left, top, width) -> None:
    """水平区切り線（細長い矩形）をアクセントカラーで描画する"""
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left, top, width, Emu(30000),  # 高さ約 0.03 インチ
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()  # 枠線なし


# ---------------------------------------------------------------------------
# 表紙スライド
# ---------------------------------------------------------------------------


def _build_title_slide(prs: Presentation, data: dict) -> None:
    """表紙（スライド 1）を生成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
    _set_bg(slide, COLOR_BG)

    w = prs.slide_width
    h = prs.slide_height

    # アクセントバー（上部）
    _add_divider(slide, Inches(0), Inches(0.6), w)

    # メインタイトル
    _add_textbox(
        slide, data.get("title", ""),
        Inches(0.8), Inches(1.2), w - Inches(1.6), Inches(1.8),
        font_size=36, bold=True, color=COLOR_TITLE, align=PP_ALIGN.CENTER,
    )

    # アクセントバー（中央）
    _add_divider(slide, Inches(2.0), Inches(3.1), w - Inches(4.0))

    # サブテキスト（日付・期間・情報源）
    body = "\n".join(data.get("content", []))
    _add_textbox(
        slide, body,
        Inches(0.8), Inches(3.4), w - Inches(1.6), Inches(2.5),
        font_size=18, color=COLOR_BODY, align=PP_ALIGN.CENTER,
    )

    # フッター
    _add_textbox(
        slide, "Space Business Weekly Report  |  Powered by Gemini API",
        Inches(0.5), h - Inches(0.6), w - Inches(1.0), Inches(0.5),
        font_size=9, color=COLOR_SOURCE, align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# 通常コンテンツスライド
# ---------------------------------------------------------------------------


def _build_content_slide(prs: Presentation, data: dict) -> None:
    """通常のコンテンツスライドを生成する"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, COLOR_BG)

    w = prs.slide_width
    h = prs.slide_height

    slide_num = data.get("slide_number", "")
    title = data.get("title", "")
    content_items = data.get("content", [])
    sources = data.get("sources", [])

    # スライド番号（右上）
    _add_textbox(
        slide, f"{slide_num} / 30",
        w - Inches(1.5), Inches(0.1), Inches(1.3), Inches(0.35),
        font_size=10, color=COLOR_NUM, align=PP_ALIGN.RIGHT,
    )

    # タイトル
    _add_textbox(
        slide, title,
        Inches(0.5), Inches(0.1), w - Inches(2.0), Inches(0.9),
        font_size=26, bold=True, color=COLOR_TITLE,
    )

    # 区切り線
    _add_divider(slide, Inches(0.5), Inches(1.05), w - Inches(1.0))

    # 本文（箇条書き）
    body_text = "\n".join(f"▶  {item}" for item in content_items)
    _add_textbox(
        slide, body_text,
        Inches(0.5), Inches(1.2), w - Inches(1.0), h - Inches(2.6),
        font_size=17, color=COLOR_BODY,
    )

    # 出典（フッター）
    if sources:
        source_str = "出典: " + " | ".join(str(s) for s in sources[:3])
        _add_textbox(
            slide, source_str,
            Inches(0.5), h - Inches(0.65), w - Inches(1.0), Inches(0.55),
            font_size=8, color=COLOR_SOURCE,
        )


# ---------------------------------------------------------------------------
# メイン
# ---------------------------------------------------------------------------


def generate() -> Path:
    """analyzed_report.json から PPTX を生成して output/ に保存する"""
    ensure_dirs()

    if not INPUT_PATH.exists():
        raise FileNotFoundError(
            f"{INPUT_PATH} が見つかりません。先に analyzer.py を実行してください。"
        )

    slides_data: list[dict] = json.loads(INPUT_PATH.read_text(encoding="utf-8"))
    logger.info("=== スライド生成開始: %d 枚 ===", len(slides_data))

    prs = Presentation()
    # 16:9 ワイドスクリーン
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides_data):
        if i == 0:
            _build_title_slide(prs, slide_data)
        else:
            _build_content_slide(prs, slide_data)
        logger.debug("  スライド %d 生成", slide_data.get("slide_number", i + 1))

    output_path = OUTPUT_DIR / f"space_report_{date.today().strftime('%Y%m%d')}.pptx"
    prs.save(str(output_path))
    logger.info("PPTX 保存完了: %s", output_path)

    return output_path


if __name__ == "__main__":
    generate()
